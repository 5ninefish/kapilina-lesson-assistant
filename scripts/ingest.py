#!/usr/bin/env python3
"""
Kaipilina Noeau — lesson ingest pipeline
Reads DOCX/PPTX/PDF, chunks, embeds with mE5-multilingual, stores in ChromaDB.
Idempotent: SHA-256 chunk IDs deduplicate on re-run.
"""

import hashlib
import json
import re
from pathlib import Path

import chromadb
from sentence_transformers import SentenceTransformer

BASE = Path.home() / "kapilinanoeau"
LESSONS_DIR = BASE / "data/lessons"
GRADE_MAP_PATH = BASE / "data/grade_map.json"
CHROMA_DIR = BASE / "data/chroma"
EMBED_MODEL = "intfloat/multilingual-e5-large"
CHUNK_CHARS = 512 * 4      # ~512 tokens
OVERLAP_CHARS = 64 * 4     # ~64 tokens
MIN_PDF_PAGE_CHARS = 100   # skip image-only pages

assert CHUNK_CHARS > OVERLAP_CHARS, "CHUNK_CHARS must exceed OVERLAP_CHARS to avoid infinite loop"


def chunk_text(text: str) -> list[str]:
    step = CHUNK_CHARS - OVERLAP_CHARS
    chunks, start = [], 0
    while start < len(text):
        chunks.append(text[start:start + CHUNK_CHARS].strip())
        start += step
    return [c for c in chunks if len(c) > 80]


def extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(path)
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    seen_cells: set[int] = set()
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if id(cell) not in seen_cells and cell.text.strip():
                    seen_cells.add(id(cell))
                    parts.append(cell.text.strip())
    return "\n".join(parts)


def extract_pptx(path: Path) -> str:
    from pptx import Presentation
    parts = []
    for slide in Presentation(path).slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


def extract_pdf(path: Path) -> str:
    import pdfplumber
    parts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            if len(text.strip()) >= MIN_PDF_PAGE_CHARS:
                parts.append(text.strip())
    return "\n".join(parts)


def parse_lesson_num(name: str) -> str | None:
    m = re.match(r"(L\d+)", name, re.IGNORECASE)
    return m.group(1).upper() if m else None


def main():
    with open(GRADE_MAP_PATH) as f:
        grade_map = json.load(f)

    print(f"Loading {EMBED_MODEL} ...")
    model = SentenceTransformer(EMBED_MODEL)

    client = chromadb.PersistentClient(path=str(CHROMA_DIR))
    collection = client.get_or_create_collection(
        name="lessons",
        metadata={"hnsw:space": "cosine"},
    )

    lessons_root = LESSONS_DIR.resolve()
    files = sorted(
        f for f in LESSONS_DIR.glob("*")
        if f.suffix.lower() in {".docx", ".pptx", ".pdf"}
        and f.resolve().is_relative_to(lessons_root)
    )
    print(f"{len(files)} lesson files found\n")

    for path in files:
        lesson_num = parse_lesson_num(path.name)
        info = grade_map.get(lesson_num, {}) if lesson_num else {}
        grades = ",".join(info.get("grades", ["unknown"]))
        title = info.get("title", path.stem)

        try:
            ext = path.suffix.lower()
            if ext == ".docx":
                text = extract_docx(path)
            elif ext == ".pptx":
                text = extract_pptx(path)
            else:
                text = extract_pdf(path)
        except Exception as e:
            print(f"  SKIP {path.name}: {e}")
            continue

        if not text.strip():
            print(f"  SKIP {path.name}: no text extracted")
            continue

        chunks = chunk_text(text)
        print(f"  {path.name}  lesson={lesson_num}  grades={grades}  {len(chunks)} chunks")

        ids, embeddings, documents, metadatas = [], [], [], []
        for i, chunk in enumerate(chunks):
            doc_id = hashlib.sha256(f"{path.name}::{i}".encode()).hexdigest()[:20]
            ids.append(doc_id)
            embeddings.append(model.encode(f"passage: {chunk}", normalize_embeddings=True).tolist())
            documents.append(chunk)
            metadatas.append({
                "filename": path.name,
                "lesson": lesson_num or "unknown",
                "grades": grades,
                "title": title,
                "chunk_idx": i,
            })

        if ids:
            collection.upsert(ids=ids, embeddings=embeddings, documents=documents, metadatas=metadatas)

    total = collection.count()
    print(f"\nDone — {total} chunks in ChromaDB at {CHROMA_DIR}")


if __name__ == "__main__":
    main()
