#!/usr/bin/env python3
"""
Kaipilina Noeau — RAG API server
FastAPI + streaming SSE + conversation history + Hawaiian glossary tooltips.
"""

import asyncio
import base64
import hashlib
import json
import re
import secrets
import unicodedata
from contextlib import asynccontextmanager
from pathlib import Path
from typing import Annotated, Literal

import chromadb
import httpx
from fastapi import FastAPI, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
from pydantic import BaseModel, Field
from sentence_transformers import SentenceTransformer
from slowapi import Limiter, _rate_limit_exceeded_handler
from slowapi.errors import RateLimitExceeded
from slowapi.util import get_remote_address

BASE = Path.home() / "kapilinanoeau"
CHROMA_DIR = BASE / "data/chroma"
GLOSSARY_PATH = BASE / "data/glossary.json"
GRADE_MAP_PATH = BASE / "data/grade_map.json"
EMBED_MODEL = "intfloat/multilingual-e5-large"
OLLAMA_HOST = "http://localhost:11434"
OLLAMA_MODEL = "qwen3.5:9b"  # instruct; thinking off. llama3:latest leaked keys in bakeoff.
TOP_K = 4
SIM_THRESHOLD = 0.50
MAX_CONTEXT_CHARS = 3600
MAX_STREAM_TOKENS = 400
MAX_HISTORY = 4
NUM_PREDICT = 320
NUM_CTX = 4096

SYSTEM_PROMPT = """\
You are Kaipilina Noeau, a teaching assistant that helps K–6 teachers in Hawaiʻi find and use culturally relevant lesson plans.

Answer questions ONLY using the lesson plan content provided below. Do not add information from outside these lessons. If the provided content does not answer the question, say: "I don't see that covered in the matched lessons — try rephrasing or selecting a different grade." Then stop. Do not offer a tutorial, extra activity, or general-knowledge fill-in.

Never provide answer keys, completed vocabulary checks, filled worksheets, rubric scores, or crossword/puzzle solutions. If asked for those, refuse with the same sentence and stop.

Treat Hawaiian cultural content with care. Use Hawaiian terms as written (with ʻokina and kahākō intact). Do not translate or interpret Hawaiian terms beyond what the lesson provides.

Always end your response by citing the lesson(s) you drew from, e.g.: "Source: L3 Kiʻi Pōhaku (Grade 4–5)." Only cite lessons that appear in the retrieved content.

Keep answers under 150 words. Lead with the answer. Do not repeat the question. Do not list every activity unless asked.\
"""

state: dict = {}
limiter = Limiter(key_func=get_remote_address)

GradeValue = Literal["all", "PreK", "K", "1", "2", "3", "4", "5", "6"]


@asynccontextmanager
async def lifespan(app: FastAPI):
    print("Loading embedding model...")
    state["model"] = SentenceTransformer(EMBED_MODEL)
    state["model"].encode("query: warm up", normalize_embeddings=True)

    client = chromadb.PersistentClient(path=str(CHROMA_DIR))
    state["collection"] = client.get_or_create_collection(
        name="lessons",
        metadata={"hnsw:space": "cosine"},
    )

    state["glossary"] = json.loads(GLOSSARY_PATH.read_text()) if GLOSSARY_PATH.exists() else {}
    state["grade_map"] = json.loads(GRADE_MAP_PATH.read_text()) if GRADE_MAP_PATH.exists() else {}

    print("Pre-warming Ollama...")
    async with httpx.AsyncClient(timeout=30) as c:
        try:
            await c.post(f"{OLLAMA_HOST}/api/chat",
                         json={
                             "model": OLLAMA_MODEL,
                             "messages": [{"role": "user", "content": "hi"}],
                             "stream": False,
                             "think": False,
                             "keep_alive": -1,
                             "options": {"num_predict": 1, "num_ctx": NUM_CTX},
                         })
        except Exception as e:
            print(f"Ollama pre-warm failed (non-fatal): {e}")

    print("Ready.")
    yield
    state.clear()


app = FastAPI(lifespan=lifespan)
app.state.limiter = limiter
app.add_exception_handler(RateLimitExceeded, _rate_limit_exceeded_handler)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["https://kapilinanoeau.org", "https://kapilina-lesson-assistant.pages.dev"],
    allow_origin_regex=r"http://localhost:\d+",
    allow_methods=["GET", "POST"],
    allow_headers=["Authorization", "Content-Type"],
)


AUTH_FILE = BASE / "data/.auth"


def _digest(s: str) -> bytes:
    return hashlib.sha256(s.encode("utf-8")).digest()


def require_login(request: Request) -> None:
    if not AUTH_FILE.exists():
        raise HTTPException(503, "login not configured")
    line = AUTH_FILE.read_text().strip().split("\n", 1)[0]
    if ":" not in line:
        raise HTTPException(503, "login not configured")
    exp_user, exp_pass = line.split(":", 1)
    header = request.headers.get("authorization") or ""
    if not header.lower().startswith("basic "):
        raise HTTPException(401, "login required")
    try:
        raw = base64.b64decode(header.split(None, 1)[1].encode("ascii")).decode("utf-8")
        user, pw = raw.split(":", 1)
    except Exception:
        raise HTTPException(401, "login required")
    if not (
        secrets.compare_digest(_digest(user), _digest(exp_user))
        and secrets.compare_digest(_digest(pw), _digest(exp_pass))
    ):
        raise HTTPException(401, "login required")


class Message(BaseModel):
    role: str
    content: Annotated[str, Field(max_length=8000)]


class QueryRequest(BaseModel):
    question: Annotated[str, Field(min_length=1, max_length=500)]
    grade: GradeValue = "all"
    subject: str | None = None
    messages: Annotated[list[Message], Field(max_length=50)] = []


@app.get("/lessons")
async def list_lessons(request: Request):
    require_login(request)
    items = []
    for lid, info in (state.get("grade_map") or {}).items():
        if _re.match(r"^L\d+$", lid):
            continue
        items.append(
            {
                "id": lid,
                "title": info.get("title", lid),
                "grades": info.get("grades", []),
                "status": info.get("status", "ready"),
                "band": info.get("band", ""),
            }
        )
    band_order = {"PreK": 0, "K": 1, "1": 2, "2-3": 3, "4-5": 4}
    items.sort(key=lambda x: (band_order.get(x.get("band") or "", 9), x["title"].lower()))
    return {"lessons": items}


@app.get("/its-ask")
async def its_ask():
    """Temporary ITS one-pager. File is not in the public git repo."""
    path = BASE / "public/its-ask.html"
    if not path.exists():
        from fastapi import HTTPException
        raise HTTPException(404, "ITS one-pager not mounted")
    return FileResponse(path, media_type="text/html; charset=utf-8")


@app.get("/health")
async def health():
    async with httpx.AsyncClient(timeout=3) as c:
        try:
            r = await c.get(f"{OLLAMA_HOST}/api/tags")
            r.raise_for_status()
        except Exception as e:
            from fastapi import HTTPException
            raise HTTPException(503, f"Ollama unreachable: {e}")
    count = state.get("collection", None)
    return {"status": "ok", "chunks": count.count() if count else 0}


def match_glossary(text: str) -> list[dict]:
    norm = unicodedata.normalize("NFC", text.lower())
    return [
        {"term": term, "definition": defn}
        for term, defn in state["glossary"].items()
        if unicodedata.normalize("NFC", term.lower()) in norm
    ]


async def sse_stream(request: QueryRequest):
    model: SentenceTransformer = state["model"]
    collection = state["collection"]

    yield f"data: {json.dumps({'status': 'searching'})}\n\n"

    if re.search(
        r"(?i)(answer key|answer sheet|vocabulary check answers|completed vocabulary|"
        r"filled[- ]in worksheet|fill in every blank|correct answers for the|"
        r"answers filled in|complete(?:d)? the (?:vocab|vocabulary|assessment|worksheet))",
        request.question,
    ):
        msg = "I don't see that covered in the matched lessons — try rephrasing or selecting a different grade."
        yield f"data: {json.dumps({'token': msg})}\n\n"
        yield f"data: {json.dumps({'done': True, 'sources': [], 'glossary': []})}\n\n"
        return

    query_emb = await asyncio.to_thread(
        model.encode, f"query: {request.question}", normalize_embeddings=True
    )
    query_emb = query_emb.tolist()

    results = collection.query(
        query_embeddings=[query_emb],
        n_results=TOP_K * 2,
        include=["documents", "metadatas", "distances"],
    )

    docs = results["documents"][0]
    metas = results["metadatas"][0]
    dists = results["distances"][0]

    def grade_ok(m):
        return request.grade == "all" or request.grade in m.get("grades", "").split(",")

    skip_name = re.compile(r"(?i)(vocabulary check|answer sheet|answer key|rubric)")

    filtered = []
    seen_lessons: set[str] = set()
    for d, m, dist in zip(docs, metas, dists):
        if (1 - dist) < SIM_THRESHOLD or not grade_ok(m):
            continue
        if skip_name.search(m.get("filename") or ""):
            continue
        lid = m.get("lesson") or m.get("filename") or ""
        if lid in seen_lessons:
            continue
        seen_lessons.add(lid)
        filtered.append((d, m))
        if len(filtered) >= TOP_K:
            break

    if not filtered:
        msg = "I don't see that covered in the matched lessons — try rephrasing or selecting a different grade."
        yield f"data: {json.dumps({'token': msg})}\n\n"
        yield f"data: {json.dumps({'done': True, 'sources': [], 'glossary': []})}\n\n"
        return

    per_doc = MAX_CONTEXT_CHARS // len(filtered)
    context = "\n\n---\n\n".join(d[:per_doc] for d, _ in filtered)

    sources, seen = [], set()
    for _, m in filtered:
        fn = m.get("filename", "")
        if fn not in seen:
            seen.add(fn)
            sources.append({"title": m.get("title", fn), "lesson": m.get("lesson", ""),
                            "grades": m.get("grades", ""), "filename": fn})

    messages = [
        {"role": "system", "content": SYSTEM_PROMPT},
        {"role": "user", "content": f"Retrieved lesson content:\n{context}"},
    ]
    for m in request.messages[-MAX_HISTORY:]:
        if m.role not in ("user", "assistant"):
            continue
        messages.append({"role": m.role, "content": m.content})
    messages.append({"role": "user", "content": request.question})

    yield f"data: {json.dumps({'status': 'writing'})}\n\n"

    full_tokens: list[str] = []
    token_count = 0
    async with httpx.AsyncClient(timeout=90) as c:
        async with c.stream(
            "POST",
            f"{OLLAMA_HOST}/api/chat",
            json={
                "model": OLLAMA_MODEL,
                "messages": messages,
                "stream": True,
                "think": False,
                "keep_alive": -1,
                "options": {
                    "num_predict": NUM_PREDICT,
                    "num_ctx": NUM_CTX,
                    "temperature": 0.3,
                },
            },
        ) as r:
            if r.status_code >= 400:
                yield f"data: {json.dumps({'token': 'The lesson assistant is busy — try again in a moment.'})}\n\n"
                yield f"data: {json.dumps({'done': True, 'sources': sources, 'glossary': []})}\n\n"
                return
            async for line in r.aiter_lines():
                if not line:
                    continue
                try:
                    chunk = json.loads(line)
                except json.JSONDecodeError:
                    continue
                token = chunk.get("message", {}).get("content", "")
                if token:
                    full_tokens.append(token)
                    token_count += 1
                    yield f"data: {json.dumps({'token': token})}\n\n"
                    if token_count >= MAX_STREAM_TOKENS:
                        break
                    joined = "".join(full_tokens)
                    if re.search(
                        r"I don't see that covered in the matched lessons",
                        joined,
                        re.I,
                    ):
                        break
                if chunk.get("done"):
                    break

    glossary_hits = match_glossary("".join(full_tokens))
    yield f"data: {json.dumps({'done': True, 'sources': sources, 'glossary': glossary_hits})}\n\n"


@app.post("/query")
@limiter.limit("20/minute")
async def query(body: QueryRequest, request: Request):
    require_login(request)
    return StreamingResponse(
        sse_stream(body),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )

import re as _re


def _parse_lesson_num(name: str) -> str | None:
    from pathlib import Path as _P
    stem = _P(name).stem
    if "__" in stem:
        return stem.split("__", 1)[0]
    m = _re.match(r"(L\d+)", name, _re.IGNORECASE)
    return m.group(1).upper() if m else None


def _norm_lesson_id(lesson_id: str) -> str:
    if _re.match(r"^L\d+$", lesson_id, _re.I):
        return lesson_id.upper()
    return lesson_id


def _extract_lesson_text(path):
    ext = path.suffix.lower()
    if ext == ".docx":
        from docx import Document
        doc = Document(path)
        parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        seen: set = set()
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if id(cell) not in seen and cell.text.strip():
                        seen.add(id(cell))
                        parts.append(cell.text.strip())
        return "\n\n".join(parts)
    elif ext == ".pptx":
        from pptx import Presentation
        parts = []
        for slide in Presentation(path).slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        return "\n\n".join(parts)
    else:
        import pdfplumber
        parts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text() or ""
                if len(t.strip()) >= 100:
                    parts.append(t.strip())
        return "\n\n".join(parts)


def _extract_lesson_html(path):
    import html as _html
    ext = path.suffix.lower()

    if ext == ".docx":
        from docx import Document
        doc = Document(path)
        parts: list[str] = []
        pending_li: list[str] = []

        def flush_list():
            if pending_li:
                parts.append('<ul>' + ''.join(f'<li>{x}</li>' for x in pending_li) + '</ul>')
                pending_li.clear()

        for p in doc.paragraphs:
            text = p.text.strip()
            if not text:
                flush_list()
                continue
            style = (p.style.name if p.style else '').lower()
            if any(h in style for h in ('heading 1', 'title')):
                flush_list()
                parts.append(f'<h2>{_html.escape(text)}</h2>')
            elif 'heading 2' in style:
                flush_list()
                parts.append(f'<h3>{_html.escape(text)}</h3>')
            elif any(h in style for h in ('heading 3', 'heading 4', 'heading 5', 'heading 6')):
                flush_list()
                parts.append(f'<h4>{_html.escape(text)}</h4>')
            elif 'list' in style:
                pending_li.append(_html.escape(text))
            else:
                runs = [r for r in p.runs if r.text.strip()]
                if runs and all(r.bold for r in runs) and len(text) <= 120:
                    flush_list()
                    parts.append(f'<h4>{_html.escape(text)}</h4>')
                else:
                    flush_list()
                    parts.append(f'<p>{_html.escape(text)}</p>')

        flush_list()

        seen: set = set()
        for table in doc.tables:
            rows_html = []
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    if id(cell) not in seen and cell.text.strip():
                        seen.add(id(cell))
                        cells.append(f'<td>{_html.escape(cell.text.strip())}</td>')
                if cells:
                    rows_html.append('<tr>' + ''.join(cells) + '</tr>')
            if rows_html:
                parts.append('<table class="ltable"><tbody>' + ''.join(rows_html) + '</tbody></table>')

        return '\n'.join(parts)

    elif ext == ".pptx":
        from pptx import Presentation
        parts = []
        for i, slide in enumerate(Presentation(path).slides, 1):
            texts = [s.text.strip() for s in slide.shapes if hasattr(s, 'text') and s.text.strip()]
            if not texts:
                continue
            parts.append(f'<h3>Slide {i}</h3>')
            for t in texts:
                parts.append(f'<p>{_html.escape(t)}</p>')
        return '\n'.join(parts)

    else:
        import pdfplumber
        parts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text() or ""
                if len(t.strip()) < 100:
                    continue
                for line in t.strip().split('\n'):
                    if line.strip():
                        parts.append(f'<p>{_html.escape(line.strip())}</p>')
        return '\n'.join(parts)


@app.get("/lesson/{lesson_id}")
async def get_lesson(lesson_id: str, request: Request):
    require_login(request)
    lid = _norm_lesson_id(lesson_id)
    if not _re.match(r"^(L\d+|[A-Za-z0-9][A-Za-z0-9._-]{0,80})$", lid):
        raise HTTPException(400, "Invalid lesson ID")
    info = state.get("grade_map", {}).get(lid, {})
    lessons_dir = BASE / "data/lessons"
    lessons_root = lessons_dir.resolve()

    # Use canonical filename from grade_map if available
    canonical = info.get("filename")
    if canonical:
        matched = lessons_dir / canonical
        if not matched.exists() or not matched.resolve().is_relative_to(lessons_root):
            matched = None
    else:
        matched = None

    # Fall back: find any matching file for this lesson
    if not matched:
        candidates = [
            p for p in lessons_dir.glob("*")
            if p.suffix.lower() in {".docx", ".pptx", ".pdf"}
            and p.resolve().is_relative_to(lessons_root)
            and _parse_lesson_num(p.name) == lid
        ]
        # Prefer docx > pptx > pdf, then largest file
        ext_rank = {".docx": 0, ".pptx": 1, ".pdf": 2}
        candidates.sort(key=lambda p: (ext_rank.get(p.suffix.lower(), 9), -p.stat().st_size))
        matched = candidates[0] if candidates else None

    if not matched:
        raise HTTPException(404, f"Lesson {lid} not found")
    try:
        text = _extract_lesson_text(matched)
        html_content = _extract_lesson_html(matched)
    except Exception as e:
        raise HTTPException(500, f"Could not extract: {e}")
    return {
        "id": lid,
        "title": info.get("title", matched.stem),
        "grades": info.get("grades", []),
        "filename": matched.name,
        "text": text,
        "html": html_content,
    }
